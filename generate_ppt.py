from pptx import Presentation
from pptx.util import Inches

prs = Presentation()

# ---------------------------
# Slide 1
# ---------------------------
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)

slide.shapes.title.text = "Engineering Automation with Agentic AI"

slide.placeholders[1].text = (
    "IIT Delhi – AI for Future Tech Leaders\n"
    "Virtual Engineering Co-worker MVP"
)

# ---------------------------
# Slide 2
# ---------------------------
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)

slide.shapes.title.text = "Problem Statement"

content = slide.placeholders[1]

content.text = (
    "Engineering teams spend significant time on:\n\n"
    "• Document review\n"
    "• Technical summarization\n"
    "• Risk validation\n"
    "• Approval workflows\n"
    "• Engineering recommendations\n\n"
    "Objective:\n"
    "Build an Agentic Engineering Decision Intelligence MVP"
)

# ---------------------------
# Slide 3
# ---------------------------
slide = prs.slides.add_slide(slide_layout)

slide.shapes.title.text = "Agentic AI Workflow"

content = slide.placeholders[1]

content.text = (
    "Document Upload\n"
    "↓\n"
    "Knowledge Agent\n"
    "↓\n"
    "Recommendation Agent\n"
    "↓\n"
    "Validation Agent\n"
    "↓\n"
    "Decision Agent\n"
    "↓\n"
    "Human Approval Agent"
)

# ---------------------------
# Slide 4
# ---------------------------
slide = prs.slides.add_slide(slide_layout)

slide.shapes.title.text = "System Architecture"

content = slide.placeholders[1]

content.text = (
    "Frontend: Streamlit\n"
    "LLM Engine: Groq\n"
    "Document Parsing: PDF / DOCX / TXT\n"
    "AI Workflow: Multi-Agent Simulation\n"
    "Output: Engineering Summary + PDF Report\n"
    "Human-in-the-loop Approval Workflow"
)

# ---------------------------
# Slide 5
# ---------------------------
slide = prs.slides.add_slide(slide_layout)

slide.shapes.title.text = "Key Features Demonstrated"

content = slide.placeholders[1]

content.text = (
    "✅ Multi-document upload\n"
    "✅ Engineering summarization\n"
    "✅ Agent workflow orchestration\n"
    "✅ Confidence scoring\n"
    "✅ Human approval workflow\n"
    "✅ AI-generated PDF reports\n"
    "✅ Streamlit cloud deployment\n"
    "✅ Groq LLM integration"
)

# ---------------------------
# Slide 6
# ---------------------------
slide = prs.slides.add_slide(slide_layout)

slide.shapes.title.text = "Future Roadmap"

content = slide.placeholders[1]

content.text = (
    "• Calculation Agents\n"
    "• RFQ Automation Agents\n"
    "• Vendor Evaluation Agents\n"
    "• Predictive Maintenance AI\n"
    "• Autonomous Engineering Workflows\n"
    "• Enterprise ERP Integration\n"
    "• Multi-Agent Orchestration"
)

# ---------------------------
# Slide 7
# ---------------------------
slide = prs.slides.add_slide(slide_layout)

slide.shapes.title.text = "Thank You"

content = slide.placeholders[1]

content.text = (
    "Engineering Automation with Agentic AI\n\n"
    "IIT Delhi – AI for Future Tech Leaders"
)

# Save PPT
prs.save("Capstone_Final_Presentation.pptx")

print("PPT generated successfully.")